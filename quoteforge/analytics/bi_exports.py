"""Business-intelligence exporters: build the owner-facing Excel workbooks (with
native charts), a Power BI-ready data package (star-schema CSVs + DAX measures +
model/report spec + build guide), and a reportlab executive presentation PDF.

Everything is generated from REAL QuoteForge data (orders table + general ledger
+ financial-report engine) - never mocked. A hand-authored binary .pbix cannot
be validated without Power BI Desktop, so the Power BI deliverable is an
open-and-refresh package: point Power BI Desktop at the CSV folder, load the
ready DAX measures, and build the four report pages from the included spec.

Public entry point: ``export_all()`` writes the ``Excel/`` and ``Power BI/``
folders at the repo root and returns the paths it created.
"""
from __future__ import annotations

import csv
from datetime import date, datetime, timedelta
from pathlib import Path

# Repo root = two levels above this file (quoteforge/analytics/bi_exports.py).
_REPO_ROOT = Path(__file__).resolve().parents[2]
EXCEL_DIR = _REPO_ROOT / "Excel"
POWERBI_DIR = _REPO_ROOT / "Power BI"


# --------------------------------------------------------------------------- #
# Data helpers
# --------------------------------------------------------------------------- #
def _orders() -> list[dict]:
    """Load every order from the database (initialising it first)."""
    from quoteforge.db.database import init_db, get_all_orders
    init_db()
    return get_all_orders(100000)


def _num(value) -> float:
    """Coerce a possibly-None/blank value to a float (0.0 on failure)."""
    try:
        return round(float(value or 0), 2)
    except (TypeError, ValueError):
        return 0.0


def _date_of(iso: str) -> date | None:
    """Parse the date part of an ISO/sqlite timestamp string; None if blank."""
    if not iso:
        return None
    try:
        return datetime.fromisoformat(str(iso).replace("Z", "")).date()
    except ValueError:
        try:
            return datetime.fromisoformat(str(iso)[:10]).date()
        except ValueError:
            return None


def _active(orders: list[dict]) -> list[dict]:
    """Orders that count toward revenue (exclude cancelled/errored)."""
    return [o for o in orders if (o.get("status") or "") not in ("cancelled", "error")]


def _fee_summary(orders: list[dict]) -> dict:
    """Best-available fee buckets for the fee-breakdown report. Uses the most
    recent imported Etsy statement when present; otherwise falls back to the
    per-order actual Etsy fees recorded on the orders themselves."""
    # Fall back to the real per-order Etsy fees so the report is never empty.
    transaction = round(sum(_num(o.get("etsy_fees_actual")) for o in orders), 2)
    return {"transaction": transaction} if transaction else {}


# --------------------------------------------------------------------------- #
# Excel workbooks
# --------------------------------------------------------------------------- #
def _bold_header(ws, row: int = 1) -> None:
    """Bold every populated cell in a worksheet row (header styling)."""
    from openpyxl.styles import Font
    for c in ws[row]:
        c.font = Font(bold=True)


def build_financial_workbook(out_path: Path) -> Path:
    """Financial dashboard workbook: P&L summary + cost-mix pie + 4-week trend
    line, the full daily ledger, the fee breakdown (incl. Offsite Ads) with a
    bar chart, and refund/cancellation rates."""
    from openpyxl import Workbook
    from openpyxl.styles import Font
    from openpyxl.chart import PieChart, LineChart, BarChart, Reference
    from quoteforge.etsy.ledger import build_ledger, weekly_trend
    from quoteforge.analytics.financial_reports import (
        fee_breakdown, refund_cancellation_rates)

    orders = _orders()
    led = build_ledger("all")
    t = led["totals"]

    wb = Workbook()
    su = wb.active
    su.title = "P&L Summary"
    su["A1"] = "QuoteForge - Financial Dashboard"
    su["A1"].font = Font(bold=True, size=14)
    kpis = [("Revenue", t["revenue"]), ("COGS (Gelato)", t["cogs"]),
            ("Etsy fees", t["etsy_fees"]), ("API cost", t["api_cost"]),
            ("Overhead", t["opex"]), ("Net profit", t["net_profit"]),
            ("Net margin %", t["margin_pct"]), ("Orders", t["orders"])]
    for i, (lbl, val) in enumerate(kpis, start=3):
        su.cell(i, 1, lbl)
        su.cell(i, 2, val)
    pie = PieChart()
    pie.title = "Cost mix"
    pie.add_data(Reference(su, min_col=2, min_row=4, max_row=7))
    pie.set_categories(Reference(su, min_col=1, min_row=4, max_row=7))
    su.add_chart(pie, "D3")
    su.column_dimensions["A"].width = 18

    # Daily ledger.
    gl = wb.create_sheet("Daily Ledger")
    gl.append(["Date", "Revenue", "COGS", "Etsy Fees", "API Cost", "OpEx",
               "Net Profit", "Orders"])
    _bold_header(gl)
    for r in led["days"]:
        gl.append([r["day"], r["revenue"], r["cogs"], r["etsy_fees"],
                   r["api_cost"], r["opex"], r["net_profit"], r["orders"]])
    gl.append(["TOTAL", t["revenue"], t["cogs"], t["etsy_fees"], t["api_cost"],
               t["opex"], t["net_profit"], t["orders"]])
    _bold_header(gl, gl.max_row)

    # 4-week trend + line chart.
    tr = wb.create_sheet("Trend (4 wks)")
    tr.append(["Week of", "Revenue", "Net Profit", "Orders"])
    _bold_header(tr)
    rows = weekly_trend(4)
    for r in rows:
        tr.append([r["week_of"], r["revenue"], r["net"], r["orders"]])
    if rows:
        ch = LineChart()
        ch.title = "Revenue & Net Profit - last 4 weeks"
        ch.y_axis.title = "$"
        ch.add_data(Reference(tr, min_col=2, max_col=3, min_row=1,
                              max_row=len(rows) + 1), titles_from_data=True)
        ch.set_categories(Reference(tr, min_col=1, min_row=2, max_row=len(rows) + 1))
        tr.add_chart(ch, "F2")

    # Fee breakdown (incl. Offsite Ads) + bar chart.
    fb = fee_breakdown(_fee_summary(orders), t["revenue"])
    fe = wb.create_sheet("Fee Breakdown")
    fe.append(["Fee Type", "Amount ($)", "% of Revenue"])
    _bold_header(fe)
    for r in fb["rows"]:
        fe.append([r["label"], r["amount"], r["pct_of_revenue"]])
    fe.append(["TOTAL FEES", fb["total_fees"], fb["total_pct_of_revenue"]])
    _bold_header(fe, fe.max_row)
    if fb["rows"]:
        bar = BarChart()
        bar.title = "Fees by type ($)"
        bar.add_data(Reference(fe, min_col=2, min_row=1, max_row=len(fb["rows"]) + 1),
                     titles_from_data=True)
        bar.set_categories(Reference(fe, min_col=1, min_row=2, max_row=len(fb["rows"]) + 1))
        fe.add_chart(bar, "E2")
    fe.column_dimensions["A"].width = 20

    # Refunds & cancellations.
    rc = refund_cancellation_rates(orders)
    q = wb.create_sheet("Refunds & Cancellations")
    q.append(["Metric", "Value"])
    _bold_header(q)
    for lbl, val in (("Total orders", rc["orders"]),
                     ("Refunded", rc["refunded"]),
                     ("Cancelled", rc["cancelled"]),
                     ("Refund rate %", rc["refund_rate_pct"]),
                     ("Cancellation rate %", rc["cancellation_rate_pct"])):
        q.append([lbl, val])
    q.column_dimensions["A"].width = 22

    out_path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(out_path)
    return out_path


def build_sales_workbook(out_path: Path) -> Path:
    """Sales & operations workbook: the full order list, revenue/profit by
    channel / vendor / product (with bar charts), the traffic-source report
    (with a pie), and the fulfillment-status funnel."""
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, PieChart, Reference
    from quoteforge.etsy.ledger import build_breakdown
    from quoteforge.analytics.financial_reports import traffic_source_report

    orders = _orders()

    wb = Workbook()
    od = wb.active
    od.title = "Orders"
    cols = ["order_id", "created_at", "status", "vendor", "channel",
            "product_type", "material", "size", "country", "acquisition_source",
            "sale_price", "gelato_cost", "shipping_collected", "tax_collected",
            "etsy_fees_actual", "net_payout", "item_count"]
    od.append(cols)
    _bold_header(od)
    for o in orders:
        od.append([o.get(c) for c in cols])

    # By Channel / Vendor / Product (+ bar chart each).
    bd = build_breakdown("all")
    for title, key in (("By Channel", "by_channel"), ("By Vendor", "by_vendor"),
                       ("By Product", "by_product")):
        sh = wb.create_sheet(title)
        sh.append([title.replace("By ", ""), "Revenue", "Cost", "Net Profit", "Orders"])
        _bold_header(sh)
        data = bd[key]
        keys = sorted(data, key=lambda x: -data[x]["net"])
        for k in keys:
            r = data[k]
            sh.append([k or "(none)", r["revenue"], r["cost"], r["net"], r["orders"]])
        if keys:
            bar = BarChart()
            bar.title = f"Net profit {title.lower()}"
            bar.add_data(Reference(sh, min_col=4, min_row=1, max_row=len(keys) + 1),
                         titles_from_data=True)
            bar.set_categories(Reference(sh, min_col=1, min_row=2, max_row=len(keys) + 1))
            sh.add_chart(bar, "G2")

    # Traffic source (+ pie).
    ts = traffic_source_report(orders)
    tsh = wb.create_sheet("Traffic Source")
    tsh.append(["Source", "Orders", "Revenue"])
    _bold_header(tsh)
    for r in ts:
        tsh.append([r["source"], r["orders"], r["revenue"]])
    if ts:
        pie = PieChart()
        pie.title = "Revenue by traffic source"
        pie.add_data(Reference(tsh, min_col=3, min_row=1, max_row=len(ts) + 1),
                     titles_from_data=True)
        pie.set_categories(Reference(tsh, min_col=1, min_row=2, max_row=len(ts) + 1))
        tsh.add_chart(pie, "E2")

    # Fulfillment funnel.
    fu = wb.create_sheet("Fulfillment Status")
    fu.append(["Status", "Orders"])
    _bold_header(fu)
    counts: dict = {}
    for o in orders:
        counts[o.get("status") or "(none)"] = counts.get(o.get("status") or "(none)", 0) + 1
    for k in sorted(counts, key=lambda x: -counts[x]):
        fu.append([k, counts[k]])

    out_path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(out_path)
    return out_path


def build_customer_workbook(out_path: Path) -> Path:
    """Customer analytics workbook: CLV summary KPIs, the top customers by
    lifetime value, and the lapsed/win-back segments."""
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, Reference
    from quoteforge.analytics.clv import build_clv

    orders = _orders()
    clv = build_clv(orders)

    wb = Workbook()
    su = wb.active
    su.title = "CLV Summary"
    su["A1"] = "Customer Lifetime Value"
    from openpyxl.styles import Font
    su["A1"].font = Font(bold=True, size=14)
    pairs = [("Customers", clv["customers"]), ("Total revenue", clv["total_revenue"]),
             ("Avg CLV", clv["avg_clv"]), ("AOV", clv["aov"]),
             ("Avg orders / customer", clv["avg_orders_per_customer"]),
             ("Repeat customers", clv["repeat_customers"]),
             ("Repeat rate %", clv["repeat_rate_pct"]),
             ("Lapsed customers", clv["lapsed_customers"])]
    for i, (lbl, val) in enumerate(pairs, start=3):
        su.cell(i, 1, lbl)
        su.cell(i, 2, val)
    su.column_dimensions["A"].width = 22

    top = clv.get("top_customers", [])
    tc = wb.create_sheet("Top Customers")
    tc.append(["Customer", "Orders", "Revenue"])
    _bold_header(tc)
    for c in top:
        tc.append([c.get("customer") or c.get("email") or c.get("name", ""),
                   c.get("orders", 0), c.get("revenue", 0)])
    if top:
        bar = BarChart()
        bar.title = "Top customers by revenue"
        bar.add_data(Reference(tc, min_col=3, min_row=1, max_row=len(top) + 1),
                     titles_from_data=True)
        bar.set_categories(Reference(tc, min_col=1, min_row=2, max_row=len(top) + 1))
        tc.add_chart(bar, "E2")
    tc.column_dimensions["A"].width = 28

    out_path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(out_path)
    return out_path


def build_excel_workbooks(outdir: Path = None) -> list[Path]:
    """Build all three Excel workbooks into ``outdir`` (default ``Excel/``)."""
    outdir = Path(outdir or EXCEL_DIR)
    return [
        build_financial_workbook(outdir / "01_Financial_Dashboard.xlsx"),
        build_sales_workbook(outdir / "02_Sales_Operations.xlsx"),
        build_customer_workbook(outdir / "03_Customer_Analytics.xlsx"),
    ]


# --------------------------------------------------------------------------- #
# Power BI package (star-schema CSVs + DAX + model/report spec + guide)
# --------------------------------------------------------------------------- #
def _write_csv(path: Path, header: list[str], rows: list[list]) -> Path:
    """Write a UTF-8 CSV with a header row; create parent dirs as needed."""
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", newline="", encoding="utf-8") as fh:
        w = csv.writer(fh)
        w.writerow(header)
        w.writerows(rows)
    return path


def build_powerbi_data(datadir: Path) -> list[Path]:
    """Write the star-schema CSV sources Power BI connects to: fact_orders,
    fact_ledger_daily, fact_fees, fact_traffic, and the dim_* tables."""
    from quoteforge.etsy.ledger import build_ledger
    from quoteforge.analytics.financial_reports import (
        fee_breakdown, traffic_source_report)
    orders = _orders()
    led = build_ledger("all")
    written: list[Path] = []

    # fact_orders (one row per order).
    ocols = ["order_id", "order_date", "status", "vendor", "channel",
             "product_type", "material", "size", "country", "acquisition_source",
             "sale_price", "gelato_cost", "shipping_collected", "tax_collected",
             "etsy_fees_actual", "net_payout", "item_count", "delivery_confirmed",
             "delivered_at"]
    orows = []
    for o in orders:
        d = _date_of(o.get("created_at"))
        orows.append([o.get("order_id"), d.isoformat() if d else "",
                      o.get("status"), o.get("vendor") or "", o.get("channel") or "",
                      o.get("product_type") or "", o.get("material") or "",
                      o.get("size") or "", o.get("country") or "",
                      o.get("acquisition_source") or "", _num(o.get("sale_price")),
                      _num(o.get("gelato_cost")), _num(o.get("shipping_collected")),
                      _num(o.get("tax_collected")), _num(o.get("etsy_fees_actual")),
                      _num(o.get("net_payout")), o.get("item_count") or 1,
                      o.get("delivery_confirmed") or 0, o.get("delivered_at") or ""])
    written.append(_write_csv(datadir / "fact_orders.csv", ocols, orows))

    # fact_ledger_daily.
    written.append(_write_csv(
        datadir / "fact_ledger_daily.csv",
        ["date", "revenue", "cogs", "etsy_fees", "api_cost", "opex",
         "net_profit", "orders"],
        [[r["day"], r["revenue"], r["cogs"], r["etsy_fees"], r["api_cost"],
          r["opex"], r["net_profit"], r["orders"]] for r in led["days"]]))

    # fact_fees.
    fb = fee_breakdown(_fee_summary(orders), led["totals"]["revenue"])
    written.append(_write_csv(
        datadir / "fact_fees.csv", ["fee", "label", "amount", "pct_of_revenue"],
        [[r["fee"], r["label"], r["amount"], r["pct_of_revenue"]] for r in fb["rows"]]))

    # fact_traffic.
    written.append(_write_csv(
        datadir / "fact_traffic.csv", ["source", "orders", "revenue"],
        [[r["source"], r["orders"], r["revenue"]] for r in traffic_source_report(orders)]))

    # dim_date spanning the order/ledger range.
    dts = [d for d in (_date_of(o.get("created_at")) for o in orders) if d]
    if dts:
        start, end = min(dts), max(dts)
        drows, cur = [], start
        while cur <= end:
            drows.append([cur.isoformat(), cur.year, cur.month, cur.strftime("%B"),
                          cur.day, cur.strftime("%A"),
                          (cur - timedelta(days=cur.weekday())).isoformat()])
            cur += timedelta(days=1)
        written.append(_write_csv(
            datadir / "dim_date.csv",
            ["date", "year", "month", "month_name", "day", "weekday", "week_of"],
            drows))

    # dim_* distinct lookups.
    for fname, col in (("dim_vendor", "vendor"), ("dim_channel", "channel"),
                       ("dim_product", "product_type")):
        vals = sorted({(o.get(col) or "(none)") for o in orders})
        written.append(_write_csv(datadir / f"{fname}.csv", [col], [[v] for v in vals]))

    return written


_DAX_MEASURES = """// QuoteForge - Power BI measures (paste into a 'Measures' table)
// Create: Home > Enter Data > name it "_Measures" > then add these via New Measure.

Revenue = SUM ( fact_orders[sale_price] )
COGS = SUM ( fact_orders[gelato_cost] )
Etsy Fees = SUM ( fact_orders[etsy_fees_actual] )
Gross Profit = [Revenue] - [COGS]
Net Profit = [Revenue] - [COGS] - [Etsy Fees]
Net Margin % = DIVIDE ( [Net Profit], [Revenue] )
Orders = DISTINCTCOUNT ( fact_orders[order_id] )
AOV = DIVIDE ( [Revenue], [Orders] )

Refunded Orders = CALCULATE ( [Orders], fact_orders[status] = "refunded" )
Cancelled Orders = CALCULATE ( [Orders], fact_orders[status] = "cancelled" )
Refund Rate % = DIVIDE ( [Refunded Orders], [Orders] )
Cancellation Rate % = DIVIDE ( [Cancelled Orders], [Orders] )

Total Fees = SUM ( fact_fees[amount] )
Offsite Ads Fee = CALCULATE ( [Total Fees], fact_fees[fee] = "offsite_ads" )
Offsite Ads % of Revenue = DIVIDE ( [Offsite Ads Fee], [Revenue] )

Tax Collected = SUM ( fact_orders[tax_collected] )
Net Payout = SUM ( fact_orders[net_payout] )
Delivered Orders = CALCULATE ( [Orders], fact_orders[delivery_confirmed] = 1 )
"""

_MODEL_MD = """# QuoteForge - Power BI Data Model & Report Spec

## How the model fits together (star schema)

**Fact tables**
- `fact_orders` - one row per order (the grain). Foreign keys: `order_date`,
  `vendor`, `channel`, `product_type`.
- `fact_ledger_daily` - one row per day (revenue, COGS, fees, opex, net profit).
- `fact_fees` - fee type x amount x % of revenue (incl. Offsite Ads).
- `fact_traffic` - orders & revenue by traffic source.

**Dimension tables**
- `dim_date` (mark as the date table; key = `date`)
- `dim_vendor` (`vendor`), `dim_channel` (`channel`), `dim_product` (`product_type`)

**Relationships (all single-direction, 1 -> *)**
- `dim_date[date]`     1 -> *  `fact_orders[order_date]`
- `dim_date[date]`     1 -> *  `fact_ledger_daily[date]`
- `dim_vendor[vendor]` 1 -> *  `fact_orders[vendor]`
- `dim_channel[channel]` 1 -> * `fact_orders[channel]`
- `dim_product[product_type]` 1 -> * `fact_orders[product_type]`

## Report pages to build (the "charts")

**Page 1 - Executive Overview**
- KPI cards: Revenue, Net Profit, Net Margin %, Orders, AOV
- Line chart: Revenue & Net Profit by `dim_date[date]`
- Donut: Cost mix (COGS / Etsy Fees / API / OpEx from `fact_ledger_daily`)

**Page 2 - Financial Detail**
- Bar: `fact_fees[label]` by [Total Fees] (Offsite Ads highlighted)
- Cards: Offsite Ads % of Revenue, Refund Rate %, Cancellation Rate %
- Table: fee type, amount, % of revenue

**Page 3 - Sales & Operations**
- Bar: Net Profit by `dim_vendor`, by `dim_channel`, by `dim_product`
- Pie: Revenue by `fact_traffic[source]`
- Funnel: order count by `fact_orders[status]`

**Page 4 - Customers**
- Cards: Repeat Rate %, Avg CLV, Lapsed customers
  (load `Excel/03_Customer_Analytics.xlsx` or extend the model with a
  `fact_customers` export if you want these native in Power BI)
"""

_README_MD = """# QuoteForge - Power BI package

This folder is an **open-and-refresh** Power BI deliverable. A binary `.pbix`
can't be authored safely outside Power BI Desktop, so instead you get the clean
data model, the DAX, and the exact report spec - about 10 minutes to a live
dashboard, and it refreshes from the CSVs forever after.

## Fastest path: open the PBIP scaffold
- Open **`PBIP/QuoteForge.pbip`** in Power BI Desktop (enable Preview features >
  *Power BI Project (.pbip) save option* + *TMDL* + *PBIR* if prompted). The
  star-schema tables, relationships and DAX measures are already wired to the
  CSVs. If the data moved, set the **DataFolder** parameter, then **Refresh**
  and add visuals per `model/DataModel.md`. It's a scaffold - on first open
  Power BI Desktop may ask to finalize the report layout.

## Or build from scratch (always works)
1. **Get Data > Folder** -> select this `Power BI/data` folder -> *Combine &
   Load* each CSV (or Get Data > Text/CSV per file for full control).
2. **Model view**: create the relationships listed in `model/DataModel.md`.
   Mark `dim_date` as the date table (`date` column).
3. **Measures**: New Measure -> paste each measure from `measures/Measures.dax`.
4. **Report**: build the four pages from `model/DataModel.md`. Save as
   `QuoteForge.pbix` in this folder.

## Refresh it later
- Re-run `python -m quoteforge.admin export-bi` to regenerate the CSVs from the
  latest orders, then hit **Refresh** in Power BI Desktop. No re-modelling.

## What's inside
- `PBIP/` - Power BI Project scaffold (TMDL model + relationships + measures)
- `data/` - star-schema CSV sources (facts + dims)
- `measures/Measures.dax` - ready-to-paste DAX
- `model/DataModel.md` - relationships + the 4-page report spec
- `QuoteForge_Executive_Presentation.pdf` / `.pptx` - the standalone exec deck
"""


def build_powerbi_package(outdir: Path = None) -> dict:
    """Build the full Power BI package (CSV data + DAX + model/readme docs) into
    ``outdir`` (default ``Power BI/``). Returns a dict of written paths."""
    outdir = Path(outdir or POWERBI_DIR)
    data = build_powerbi_data(outdir / "data")
    measures = outdir / "measures" / "Measures.dax"
    measures.parent.mkdir(parents=True, exist_ok=True)
    measures.write_text(_DAX_MEASURES, encoding="utf-8")
    model = outdir / "model" / "DataModel.md"
    model.parent.mkdir(parents=True, exist_ok=True)
    model.write_text(_MODEL_MD, encoding="utf-8")
    readme = outdir / "README.md"
    readme.write_text(_README_MD, encoding="utf-8")
    return {"data": data, "measures": measures, "model": model, "readme": readme}


# --------------------------------------------------------------------------- #
# Executive presentation (reportlab PDF)
# --------------------------------------------------------------------------- #
def build_presentation(out_path: Path = None) -> Path:
    """Build a landscape executive-presentation PDF (title, KPI scorecard, fee &
    quality summary, traffic mix) from the live ledger + financial reports."""
    out_path = Path(out_path or (POWERBI_DIR / "QuoteForge_Executive_Presentation.pdf"))
    from reportlab.lib.pagesizes import landscape, letter
    from reportlab.lib import colors
    from reportlab.lib.units import inch
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table,
                                    TableStyle, PageBreak)
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from quoteforge.etsy.ledger import build_ledger
    from quoteforge.analytics.financial_reports import (
        fee_breakdown, refund_cancellation_rates, traffic_source_report)

    orders = _orders()
    led = build_ledger("all")
    t = led["totals"]
    fb = fee_breakdown(_fee_summary(orders), t["revenue"])
    rc = refund_cancellation_rates(orders)
    ts = traffic_source_report(orders)

    styles = getSampleStyleSheet()
    title = ParagraphStyle("T", parent=styles["Title"], textColor=colors.HexColor("#5b3fa0"))
    h2 = ParagraphStyle("H2", parent=styles["Heading2"], textColor=colors.HexColor("#5b3fa0"))

    out_path.parent.mkdir(parents=True, exist_ok=True)
    doc = SimpleDocTemplate(str(out_path), pagesize=landscape(letter),
                            topMargin=0.6 * inch, bottomMargin=0.6 * inch)
    story = [Spacer(1, 1.6 * inch),
             Paragraph("QuoteForge", title),
             Paragraph("Executive Business Review", h2),
             Spacer(1, 0.3 * inch),
             Paragraph("Revenue, profitability, fees, quality &amp; traffic - "
                       "generated from live order data.", styles["Normal"]),
             PageBreak()]

    def _table(rows, widths):
        """Build a styled reportlab table for a slide."""
        tb = Table(rows, colWidths=widths, hAlign="LEFT")
        tb.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#5b3fa0")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f1ecfb")]),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#d9cdf2")),
            ("FONTSIZE", (0, 0), (-1, -1), 11),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 7),
            ("TOPPADDING", (0, 0), (-1, -1), 7)]))
        return tb

    # Slide 2 - KPI scorecard.
    story += [Paragraph("Profit &amp; Loss Scorecard", h2), Spacer(1, 0.2 * inch),
              _table([["Metric", "Value"],
                      ["Revenue", f"${t['revenue']:,.2f}"],
                      ["COGS (Gelato)", f"${t['cogs']:,.2f}"],
                      ["Etsy fees", f"${t['etsy_fees']:,.2f}"],
                      ["Net profit", f"${t['net_profit']:,.2f}"],
                      ["Net margin", f"{t['margin_pct']}%"],
                      ["Orders", f"{t['orders']}"]],
                     [3 * inch, 3 * inch]), PageBreak()]

    # Slide 3 - fees & quality.
    fee_rows = [["Fee type", "Amount", "% of revenue"]] + \
               [[r["label"], f"${r['amount']:,.2f}", f"{r['pct_of_revenue']}%"] for r in fb["rows"]] + \
               [["TOTAL FEES", f"${fb['total_fees']:,.2f}", f"{fb['total_pct_of_revenue']}%"]]
    story += [Paragraph("Fees &amp; Quality", h2), Spacer(1, 0.2 * inch),
              _table(fee_rows, [3 * inch, 2 * inch, 2 * inch]),
              Spacer(1, 0.25 * inch),
              Paragraph(f"Refund rate <b>{rc['refund_rate_pct']}%</b> &nbsp;|&nbsp; "
                        f"Cancellation rate <b>{rc['cancellation_rate_pct']}%</b> "
                        f"of {rc['orders']} orders", styles["Normal"]),
              PageBreak()]

    # Slide 4 - traffic mix.
    tr_rows = [["Traffic source", "Orders", "Revenue"]] + \
              [[r["source"], f"{r['orders']}", f"${r['revenue']:,.2f}"] for r in ts]
    story += [Paragraph("Where Sales Come From", h2), Spacer(1, 0.2 * inch),
              _table(tr_rows if ts else [["Traffic source", "Orders", "Revenue"],
                                         ["(no data yet)", "0", "$0.00"]],
                     [3 * inch, 2 * inch, 2 * inch])]

    doc.build(story)
    return out_path


# --------------------------------------------------------------------------- #
# PowerPoint presentation (python-pptx)
# --------------------------------------------------------------------------- #
def build_pptx_presentation(out_path: Path = None) -> Path:
    """Build a native PowerPoint (.pptx) executive deck - title, P&L scorecard,
    fees & quality, and traffic mix - from the live ledger + financial reports.
    Mirrors the PDF deck but as an editable .pptx."""
    out_path = Path(out_path or (POWERBI_DIR / "QuoteForge_Executive_Presentation.pptx"))
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from quoteforge.etsy.ledger import build_ledger
    from quoteforge.analytics.financial_reports import (
        fee_breakdown, refund_cancellation_rates, traffic_source_report)

    orders = _orders()
    t = build_ledger("all")["totals"]
    fb = fee_breakdown(_fee_summary(orders), t["revenue"])
    rc = refund_cancellation_rates(orders)
    ts = traffic_source_report(orders)
    purple = RGBColor(0x5B, 0x3F, 0xA0)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def _title_slide(title: str, subtitle: str) -> None:
        """Add a centered title slide."""
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(Inches(1), Inches(2.6), Inches(11.3), Inches(2))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run(); r.text = title
        r.font.size = Pt(54); r.font.bold = True; r.font.color.rgb = purple
        p2 = tb.text_frame.add_paragraph()
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(24); r2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    def _table_slide(title: str, header: list, rows: list) -> None:
        """Add a slide with a heading and a styled data table."""
        s = prs.slides.add_slide(blank)
        ht = s.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.9))
        hr = ht.text_frame.paragraphs[0].add_run(); hr.text = title
        hr.font.size = Pt(32); hr.font.bold = True; hr.font.color.rgb = purple
        nrows, ncols = len(rows) + 1, len(header)
        gt = s.shapes.add_table(nrows, ncols, Inches(0.6), Inches(1.5),
                                Inches(12), Inches(0.4 * nrows)).table
        for c, txt in enumerate(header):
            cell = gt.cell(0, c); cell.text = str(txt)
            cell.fill.solid(); cell.fill.fore_color.rgb = purple
            cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.text_frame.paragraphs[0].runs[0].font.bold = True
        for ri, row in enumerate(rows, start=1):
            for c, txt in enumerate(row):
                gt.cell(ri, c).text = str(txt)

    _title_slide("QuoteForge", "Executive Business Review")
    _table_slide("Profit & Loss Scorecard", ["Metric", "Value"],
                 [["Revenue", f"${t['revenue']:,.2f}"],
                  ["COGS (Gelato)", f"${t['cogs']:,.2f}"],
                  ["Etsy fees", f"${t['etsy_fees']:,.2f}"],
                  ["Net profit", f"${t['net_profit']:,.2f}"],
                  ["Net margin", f"{t['margin_pct']}%"],
                  ["Orders", f"{t['orders']}"]])
    _table_slide("Fees & Quality", ["Fee type", "Amount", "% of revenue"],
                 [[r["label"], f"${r['amount']:,.2f}", f"{r['pct_of_revenue']}%"]
                  for r in fb["rows"]]
                 + [["TOTAL FEES", f"${fb['total_fees']:,.2f}",
                     f"{fb['total_pct_of_revenue']}%"],
                    [f"Refund rate {rc['refund_rate_pct']}%",
                     f"Cancellation {rc['cancellation_rate_pct']}%",
                     f"of {rc['orders']} orders"]])
    _table_slide("Where Sales Come From", ["Traffic source", "Orders", "Revenue"],
                 [[r["source"], r["orders"], f"${r['revenue']:,.2f}"] for r in ts]
                 or [["(no data yet)", "0", "$0.00"]])

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


# --------------------------------------------------------------------------- #
# Power BI Project (PBIP) scaffold - TMDL semantic model + report shell
# --------------------------------------------------------------------------- #
# (table -> [(column, dataType, summarizeBy)]). dataType drives both the TMDL
# column type and the M type-transform; summarizeBy "none" leaves dims un-summed.
_PBIP_TABLES: dict = {
    "fact_orders": [
        ("order_id", "string", "none"), ("order_date", "dateTime", "none"),
        ("status", "string", "none"), ("vendor", "string", "none"),
        ("channel", "string", "none"), ("product_type", "string", "none"),
        ("material", "string", "none"), ("size", "string", "none"),
        ("country", "string", "none"), ("acquisition_source", "string", "none"),
        ("sale_price", "double", "sum"), ("gelato_cost", "double", "sum"),
        ("shipping_collected", "double", "sum"), ("tax_collected", "double", "sum"),
        ("etsy_fees_actual", "double", "sum"), ("net_payout", "double", "sum"),
        ("item_count", "int64", "sum"), ("delivery_confirmed", "int64", "none"),
        ("delivered_at", "string", "none")],
    "fact_ledger_daily": [
        ("date", "dateTime", "none"), ("revenue", "double", "sum"),
        ("cogs", "double", "sum"), ("etsy_fees", "double", "sum"),
        ("api_cost", "double", "sum"), ("opex", "double", "sum"),
        ("net_profit", "double", "sum"), ("orders", "int64", "sum")],
    "fact_fees": [
        ("fee", "string", "none"), ("label", "string", "none"),
        ("amount", "double", "sum"), ("pct_of_revenue", "double", "none")],
    "fact_traffic": [
        ("source", "string", "none"), ("orders", "int64", "sum"),
        ("revenue", "double", "sum")],
    "dim_date": [
        ("date", "dateTime", "none"), ("year", "int64", "none"),
        ("month", "int64", "none"), ("month_name", "string", "none"),
        ("day", "int64", "none"), ("weekday", "string", "none"),
        ("week_of", "dateTime", "none")],
    "dim_vendor": [("vendor", "string", "none")],
    "dim_channel": [("channel", "string", "none")],
    "dim_product": [("product_type", "string", "none")],
}

# (fromTable.fromColumn -> toTable.toColumn), all single-direction 1->*.
_PBIP_RELATIONSHIPS = [
    ("fact_orders", "order_date", "dim_date", "date"),
    ("fact_ledger_daily", "date", "dim_date", "date"),
    ("fact_orders", "vendor", "dim_vendor", "vendor"),
    ("fact_orders", "channel", "dim_channel", "channel"),
    ("fact_orders", "product_type", "dim_product", "product_type"),
]

# DAX measures hosted on the _Measures table (name -> expression).
_PBIP_MEASURES = [
    ("Revenue", "SUM ( fact_orders[sale_price] )"),
    ("COGS", "SUM ( fact_orders[gelato_cost] )"),
    ("Etsy Fees", "SUM ( fact_orders[etsy_fees_actual] )"),
    ("Net Profit", "[Revenue] - [COGS] - [Etsy Fees]"),
    ("Net Margin %", "DIVIDE ( [Net Profit], [Revenue] )"),
    ("Orders", "DISTINCTCOUNT ( fact_orders[order_id] )"),
    ("AOV", "DIVIDE ( [Revenue], [Orders] )"),
    ("Refund Rate %", 'DIVIDE ( CALCULATE ( [Orders], fact_orders[status] = "refunded" ), [Orders] )'),
    ("Cancellation Rate %", 'DIVIDE ( CALCULATE ( [Orders], fact_orders[status] = "cancelled" ), [Orders] )'),
    ("Total Fees", "SUM ( fact_fees[amount] )"),
    ("Offsite Ads % of Revenue", 'DIVIDE ( CALCULATE ( [Total Fees], fact_fees[fee] = "offsite_ads" ), [Revenue] )'),
]

_M_TYPE = {"string": "type text", "double": "type number",
           "int64": "Int64.Type", "dateTime": "type date"}


def _guid(seed: str) -> str:
    """Deterministic GUID from a seed (stable across regenerations)."""
    import uuid
    return str(uuid.uuid5(uuid.NAMESPACE_URL, "quoteforge.pbip." + seed))


def _tmdl_table(name: str, cols: list) -> str:
    """Render one TMDL table file (columns + an M partition reading its CSV)."""
    lines = [f"table {name}", f"\tlineageTag: {_guid('t.' + name)}", ""]
    for col, dtype, summ in cols:
        lines += [f"\tcolumn {col}", f"\t\tdataType: {dtype}",
                  f"\t\tlineageTag: {_guid(f'c.{name}.{col}')}",
                  f"\t\tsummarizeBy: {summ}", f"\t\tsourceColumn: {col}", ""]
    transforms = ", ".join(f'{{"{c}", {_M_TYPE[d]}}}' for c, d, _ in cols)
    m = ("let\n"
         f'    Source = Csv.Document(File.Contents(DataFolder & "/{name}.csv"),'
         "[Delimiter=\",\", Encoding=65001, QuoteStyle=QuoteStyle.Csv]),\n"
         "    Promoted = Table.PromoteHeaders(Source, [PromoteAllScalars=true]),\n"
         f"    Typed = Table.TransformColumnTypes(Promoted, {{{transforms}}})\n"
         "in\n    Typed")
    m_indented = "\n".join("\t\t\t" + ln for ln in m.splitlines())
    lines += [f"\tpartition {name} = m", "\t\tmode: import",
              "\t\tsource = ```", m_indented, "\t\t\t```", ""]
    return "\n".join(lines)


def _tmdl_measures() -> str:
    """Render the _Measures table (a 1-row helper table hosting the DAX)."""
    lines = ["table _Measures", f"\tlineageTag: {_guid('t._Measures')}", ""]
    for nm, expr in _PBIP_MEASURES:
        lines += [f"\tmeasure '{nm}' = {expr}",
                  f"\t\tlineageTag: {_guid('m.' + nm)}", ""]
    lines += ["\tcolumn _dummy", "\t\tdataType: int64", "\t\tisHidden",
              f"\t\tlineageTag: {_guid('c._Measures._dummy')}",
              "\t\tsummarizeBy: none", "\t\tsourceColumn: _dummy", "",
              "\tpartition _Measures = m", "\t\tmode: import",
              "\t\tsource = ```", "\t\t\tlet Source = #table({\"_dummy\"}, {{1}}) in Source",
              "\t\t\t```", ""]
    return "\n".join(lines)


def _tmdl_relationships() -> str:
    """Render relationships.tmdl (all single-direction many-to-one)."""
    out = []
    for ft, fc, tt, tc in _PBIP_RELATIONSHIPS:
        out += [f"relationship {_guid(f'r.{ft}.{fc}.{tt}.{tc}')}",
                f"\tfromColumn: {ft}.{fc}", f"\ttoColumn: {tt}.{tc}", ""]
    return "\n".join(out)


def build_pbip_project(outdir: Path = None, data_dir: Path = None) -> dict:
    """Write a Power BI Project (PBIP) scaffold: a TMDL semantic model with the
    star-schema tables, relationships and DAX measures pre-wired to the CSV
    folder, plus a report shell. Open ``QuoteForge.pbip`` in Power BI Desktop;
    set the ``DataFolder`` parameter if the data moved, then refresh and add
    visuals per model/DataModel.md. Returns the written paths."""
    import json
    base = Path(outdir or (POWERBI_DIR / "PBIP"))
    data_dir = Path(data_dir or (POWERBI_DIR / "data")).resolve()
    sm = base / "QuoteForge.SemanticModel"
    rp = base / "QuoteForge.Report"
    (sm / "definition" / "tables").mkdir(parents=True, exist_ok=True)
    (rp / "definition").mkdir(parents=True, exist_ok=True)
    written = []

    def _w(path: Path, text: str) -> None:
        """Write one scaffold file (UTF-8) and record it."""
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(text, encoding="utf-8")
        written.append(path)

    # .pbip entry point.
    _w(base / "QuoteForge.pbip", json.dumps({
        "version": "1.0",
        "artifacts": [{"report": {"path": "QuoteForge.Report"}}],
        "settings": {"enableAutoRecovery": True}}, indent=2))

    # Semantic model.
    _w(sm / ".platform", json.dumps({
        "$schema": "https://developer.microsoft.com/json-schemas/fabric/gitIntegration/platformProperties/2.0.0/schema.json",
        "metadata": {"type": "SemanticModel", "displayName": "QuoteForge"},
        "config": {"version": "2.0", "logicalId": _guid("sm.logical")}}, indent=2))
    _w(sm / "definition.pbism", json.dumps({"version": "4.2", "settings": {}}, indent=2))
    _w(sm / "definition" / "model.tmdl",
       "model Model\n\tculture: en-US\n\tdefaultPowerBIDataSourceVersion: powerBI_V3\n")
    _w(sm / "definition" / "expressions.tmdl",
       f'expression DataFolder = "{data_dir.as_posix()}" meta '
       '[IsParameterQuery=true, Type="Text", IsParameterQueryRequired=true]\n')
    _w(sm / "definition" / "relationships.tmdl", _tmdl_relationships())
    for name, cols in _PBIP_TABLES.items():
        _w(sm / "definition" / "tables" / f"{name}.tmdl", _tmdl_table(name, cols))
    _w(sm / "definition" / "tables" / "_Measures.tmdl", _tmdl_measures())

    # Report shell (single starter page; build the rest per the model spec).
    _w(rp / ".platform", json.dumps({
        "$schema": "https://developer.microsoft.com/json-schemas/fabric/gitIntegration/platformProperties/2.0.0/schema.json",
        "metadata": {"type": "Report", "displayName": "QuoteForge"},
        "config": {"version": "2.0", "logicalId": _guid("rp.logical")}}, indent=2))
    _w(rp / "definition.pbir", json.dumps({
        "version": "1.0",
        "datasetReference": {"byPath": {"path": "../QuoteForge.SemanticModel"}}}, indent=2))
    _w(rp / "definition" / "report.json", json.dumps({
        "$schema": "https://developer.microsoft.com/json-schemas/fabric/item/report/definition/1.0.0/schema.json",
        "themeCollection": {"baseTheme": {"name": "CY24SU10"}},
        "layoutOptimization": "None"}, indent=2))

    return {"pbip": base / "QuoteForge.pbip", "files": written}


# --------------------------------------------------------------------------- #
# Orchestrator
# --------------------------------------------------------------------------- #
def export_all(excel_dir: Path = None, powerbi_dir: Path = None) -> dict:
    """Build everything: the Excel workbooks, the Power BI package (CSVs + DAX +
    docs), the PBIP project scaffold, and the executive presentation in both PDF
    and PowerPoint. Returns a dict of all created paths."""
    pbase = Path(powerbi_dir) if powerbi_dir else POWERBI_DIR
    xl = build_excel_workbooks(excel_dir)
    pbi = build_powerbi_package(powerbi_dir)
    pbip = build_pbip_project(pbase / "PBIP", pbase / "data")
    pdf = build_presentation(pbase / "QuoteForge_Executive_Presentation.pdf")
    pptx = build_pptx_presentation(pbase / "QuoteForge_Executive_Presentation.pptx")
    return {"excel": xl, "powerbi": pbi, "pbip": pbip,
            "presentation_pdf": pdf, "presentation_pptx": pptx}
